from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

def main() -> None:
    out = Path(__file__).parent / "sample.pptx"
    pres = Presentation()
    layout = pres.slide_layouts[5]
    s1 = pres.slides.add_slide(layout)
    s1.shapes.title.text = "Slide One"
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    tb.text_frame.text = "Body of slide one."
    s1.notes_slide.notes_text_frame.text = "Speaker note one."
    s2 = pres.slides.add_slide(layout)
    s2.shapes.title.text = "Slide Two"
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
    tb2.text_frame.text = "Body of slide two."
    pres.save(str(out))

if __name__ == "__main__":
    main()
