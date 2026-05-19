from __future__ import annotations

import io

from pptx import Presentation

from core.exceptions import AppError, ErrorCode
from core.ingestion.parsers import register
from core.ingestion.types import ParsedDocument, ParsedSection


class PptxParser:
    kind = "pptx"

    def parse_bytes(self, data: bytes, *, source_hint: str | None = None) -> ParsedDocument:
        try:
            pres = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise AppError(
                ErrorCode.INGESTION_PARSE_FAILED, "pptx parse failed", detail=str(exc)
            ) from exc

        sections: list[ParsedSection] = []
        for slide_index, slide in enumerate(pres.slides):
            title = ""
            body_parts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape == slide.shapes.title:
                    title = text
                else:
                    body_parts.append(text)
            if slide.has_notes_slide:
                note = slide.notes_slide.notes_text_frame.text.strip()
                if note:
                    body_parts.append(f"[notes] {note}")
            body = "\n".join(body_parts).strip()
            if not body and not title:
                continue
            sections.append(
                ParsedSection(
                    text=body or title,
                    page=slide_index + 1,
                    heading_path=[title] if title else [],
                    ord=slide_index,
                )
            )

        if not sections:
            raise AppError(
                ErrorCode.INGESTION_PARSE_FAILED, "no text in pptx"
            )
        return ParsedDocument(title=source_hint or "deck.pptx", sections=sections)


register(PptxParser())
